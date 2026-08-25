"""Tests for grep search tools."""

from __future__ import annotations

import asyncio
import os
import re
import threading
import time
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import AsyncMock, MagicMock

import pytest

from nanobot.agent.loop import AgentLoop
from nanobot.agent.subagent import SubagentManager, SubagentStatus
from nanobot.agent.tools.search import FindFilesTool, GrepTool
from nanobot.agent.tools.web import WebSearchTool
from nanobot.bus.queue import MessageBus
from nanobot.config.schema import WebSearchConfig
from nanobot.providers.base import GenerationSettings
from nanobot.security.workspace_access import (
    bind_workspace_scope,
    default_workspace_scope,
    reset_workspace_scope,
)
from nanobot.utils.llm_runtime import LLMRuntime


@pytest.mark.asyncio
async def test_web_search_tool_refreshes_dynamic_config_loader(monkeypatch) -> None:
    tool = WebSearchTool(
        config=WebSearchConfig(provider="brave"),
        config_loader=lambda: WebSearchConfig(provider="duckduckgo", max_results=3),
    )

    async def fake_duckduckgo(self, query: str, n: int) -> str:
        return f"{self.config.provider}:{query}:{n}"

    monkeypatch.setattr(WebSearchTool, "_search_duckduckgo", fake_duckduckgo)

    assert await tool.execute("nanobot") == "duckduckgo:nanobot:3"


@pytest.mark.asyncio
async def test_find_files_filters_by_query_glob_and_type(tmp_path: Path) -> None:
    (tmp_path / "src").mkdir()
    (tmp_path / "src" / "settings_view.tsx").write_text("export {}\n", encoding="utf-8")
    (tmp_path / "src" / "settings_api.py").write_text("pass\n", encoding="utf-8")
    (tmp_path / "README.md").write_text("settings\n", encoding="utf-8")

    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        path=".",
        query="settings",
        glob="src/**",
        type="ts",
    )

    assert result.splitlines() == ["src/settings_view.tsx"]


@pytest.mark.asyncio
async def test_find_files_can_include_directories(tmp_path: Path) -> None:
    (tmp_path / "src" / "settings").mkdir(parents=True)
    (tmp_path / "src" / "settings" / "index.ts").write_text("export {}\n", encoding="utf-8")

    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(path="src", query="settings", include_dirs=True)

    assert "src/settings/" in result.splitlines()
    assert "src/settings/index.ts" in result.splitlines()


@pytest.mark.asyncio
async def test_find_files_supports_modified_sort_and_pagination(tmp_path: Path) -> None:
    (tmp_path / "src").mkdir()
    for idx, name in enumerate(("a.py", "b.py", "c.py"), start=1):
        file_path = tmp_path / "src" / name
        file_path.write_text("pass\n", encoding="utf-8")
        os.utime(file_path, (idx, idx))

    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        path="src",
        type="py",
        sort="modified",
        head_limit=1,
        offset=1,
    )

    assert result.splitlines()[0] == "src/b.py"
    assert "pagination: limit=1, offset=1" in result


@pytest.mark.asyncio
async def test_find_files_rejects_paths_outside_workspace(tmp_path: Path) -> None:
    outside = tmp_path.parent / "outside-find-files.txt"
    outside.write_text("secret\n", encoding="utf-8")

    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(path=str(outside))

    assert result.startswith("Error:")


@pytest.mark.asyncio
async def test_find_files_worker_preserves_current_workspace_scope(tmp_path: Path) -> None:
    project = tmp_path / "project"
    project.mkdir()
    (project / "inside.txt").write_text("ok\n", encoding="utf-8")
    (tmp_path / "outside.txt").write_text("nope\n", encoding="utf-8")
    tool = FindFilesTool(workspace=tmp_path, restrict_to_workspace=False)
    token = bind_workspace_scope(default_workspace_scope(project, True))
    try:
        result = await tool.execute(path=".")
    finally:
        reset_workspace_scope(token)

    assert result == "inside.txt"


@pytest.mark.asyncio
async def test_find_files_scan_keeps_event_loop_responsive(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    target = tmp_path / "match.txt"
    target.write_text("ok\n", encoding="utf-8")
    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)
    original_iter_paths = tool._iter_paths
    started = threading.Event()
    release = threading.Event()

    def blocking_iter_paths(root: Path, *, include_dirs: bool, budget):
        started.set()
        if not release.wait(timeout=1):
            raise TimeoutError("test did not release find_files traversal")
        yield from original_iter_paths(root, include_dirs=include_dirs, budget=budget)

    monkeypatch.setattr(tool, "_iter_paths", blocking_iter_paths)
    task = asyncio.create_task(tool.execute(path="."))
    try:
        assert await asyncio.to_thread(started.wait, 0.5)
        for _ in range(3):
            await asyncio.sleep(0.01)
        assert not task.done()
    finally:
        release.set()

    assert await asyncio.wait_for(task, timeout=0.5) == "match.txt"


@pytest.mark.asyncio
async def test_find_files_cancellation_stops_worker_scan(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)
    started = threading.Event()
    stopped = threading.Event()

    def cancellable_iter_paths(root: Path, *, include_dirs: bool, budget):
        del root, include_dirs
        started.set()
        if not budget.cancelled.wait(timeout=1):
            raise TimeoutError("find_files worker did not receive cancellation")
        stopped.set()
        if False:
            yield

    monkeypatch.setattr(tool, "_iter_paths", cancellable_iter_paths)
    task = asyncio.create_task(tool.execute(path="."))
    assert await asyncio.to_thread(started.wait, 0.5)

    task.cancel()
    with pytest.raises(asyncio.CancelledError):
        await asyncio.wait_for(task, timeout=0.2)

    assert await asyncio.to_thread(stopped.wait, 0.5)


@pytest.mark.asyncio
async def test_find_files_path_limit_stops_after_pagination_lookahead(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    paths = [tmp_path / name for name in ("a.txt", "b.txt", "c.txt")]
    for path in paths:
        path.write_text("ok\n", encoding="utf-8")
    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)
    visited: list[str] = []

    def ordered_iter_paths(root: Path, *, include_dirs: bool, budget):
        del include_dirs
        for path in paths:
            budget.checkpoint()
            visited.append(path.name)
            yield tool._entry(path, root, is_dir=False)

    monkeypatch.setattr(tool, "_iter_paths", ordered_iter_paths)
    result = await tool.execute(path=".", head_limit=1)

    assert result.splitlines()[0] == "a.txt"
    assert "pagination: limit=1, offset=0" in result
    assert visited == ["a.txt", "b.txt"]


@pytest.mark.asyncio
async def test_find_files_path_budget_counts_directories(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    (tmp_path / "one").mkdir()
    (tmp_path / "two").mkdir()
    monkeypatch.setattr(FindFilesTool, "_MAX_SCAN_PATHS", 1)
    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)

    result = await tool.execute(path=".")

    assert result.startswith("Error: find_files scan exceeded 1 paths")


@pytest.mark.asyncio
async def test_find_files_enforces_time_budget(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    (tmp_path / "match.txt").write_text("ok\n", encoding="utf-8")
    monkeypatch.setattr(FindFilesTool, "_MAX_SCAN_SECONDS", 0.0)
    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)

    result = await tool.execute(path=".")

    assert result.startswith("Error: find_files scan exceeded 0 seconds")


@pytest.mark.asyncio
async def test_find_files_path_sort_matches_existing_lexicographic_contract(
    tmp_path: Path,
) -> None:
    (tmp_path / "a").mkdir()
    (tmp_path / "a" / "inside.txt").write_text("ok\n", encoding="utf-8")
    (tmp_path / "a+").write_text("ok\n", encoding="utf-8")
    (tmp_path / "a.py").write_text("ok\n", encoding="utf-8")
    tool = FindFilesTool(workspace=tmp_path, allowed_dir=tmp_path)

    result = await tool.execute(path=".", head_limit=0)

    assert result.splitlines() == ["a+", "a.py", "a/inside.txt"]


@pytest.mark.asyncio
async def test_grep_respects_glob_filter_and_context(tmp_path: Path) -> None:
    (tmp_path / "src").mkdir()
    (tmp_path / "src" / "main.py").write_text(
        "alpha\nbeta\nmatch_here\ngamma\n",
        encoding="utf-8",
    )
    (tmp_path / "README.md").write_text("match_here\n", encoding="utf-8")

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="match_here",
        path=".",
        glob="*.py",
        output_mode="content",
        context_before=1,
        context_after=1,
    )

    assert "src/main.py:3" in result
    assert "  2| beta" in result
    assert "> 3| match_here" in result
    assert "  4| gamma" in result
    assert "README.md" not in result


@pytest.mark.asyncio
async def test_grep_defaults_to_match_context(tmp_path: Path) -> None:
    (tmp_path / "src").mkdir()
    (tmp_path / "src" / "main.py").write_text(
        "\n".join(f"line {line}" for line in range(1, 6))
        + "\nmatch_here\n"
        + "\n".join(f"line {line}" for line in range(7, 13)),
        encoding="utf-8",
    )

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="match_here",
        path="src",
    )

    assert "src/main.py:6" in result
    assert "  1| line 1" in result
    assert "> 6| match_here" in result
    assert "  11| line 11" in result
    assert "line 12" not in result


@pytest.mark.asyncio
async def test_grep_searches_xlsx_with_sheet_cell_locator(tmp_path: Path) -> None:
    from openpyxl import Workbook

    workbook_path = tmp_path / "people.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "People"
    sheet.append(["Name", "Role"])
    sheet.append(["Ada", "Engineer"])
    workbook.save(workbook_path)
    workbook.close()

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="Engineer",
        path="people.xlsx",
        fixed_strings=True,
    )

    assert "people.xlsx:3" in result
    assert "sheet='People',row=2,cell=B2" in result
    assert "Ada\tEngineer" in result


@pytest.mark.asyncio
async def test_grep_searches_docx_with_paragraph_locator(tmp_path: Path) -> None:
    from docx import Document

    document_path = tmp_path / "notes.docx"
    document = Document()
    document.add_paragraph("Introduction")
    document.add_paragraph("late needle")
    document.save(document_path)

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="late needle",
        path="notes.docx",
        fixed_strings=True,
    )

    assert "notes.docx:3 [paragraph=2]" in result
    assert "late needle" in result


@pytest.mark.asyncio
async def test_grep_searches_pptx_with_slide_locator(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    )
    textbox.text_frame.text = "slide needle"
    presentation.save(presentation_path)

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="slide needle",
        path="deck.pptx",
        fixed_strings=True,
    )

    assert "deck.pptx:2 [slide=1,line=1]" in result
    assert "slide needle" in result


@pytest.mark.asyncio
async def test_grep_searches_pdf_with_page_locator(tmp_path: Path) -> None:
    import fitz

    pdf_path = tmp_path / "notes.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "pdf needle")
    document.save(pdf_path)
    document.close()

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="pdf needle",
        path="notes.pdf",
        fixed_strings=True,
    )

    assert "notes.pdf:2 [page=1,line=1]" in result
    assert "pdf needle" in result


@pytest.mark.asyncio
async def test_grep_searches_xlsx_beyond_attachment_preview_limit(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from openpyxl import Workbook

    from nanobot.utils import document as document_utils

    workbook_path = tmp_path / "long.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    for row in range(1, 20):
        sheet.append([f"ordinary-row-{row}"])
    sheet.append(["late-needle"])
    workbook.save(workbook_path)
    workbook.close()
    monkeypatch.setattr(document_utils, "_MAX_TEXT_LENGTH", 50)

    preview = document_utils.extract_text(workbook_path)
    assert preview is not None
    assert "late-needle" not in preview

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="late-needle",
        path="long.xlsx",
        fixed_strings=True,
    )

    assert "late-needle" in result
    assert "sheet='Data',row=20,cell=A20" in result
    assert "No matches found" not in result


@pytest.mark.asyncio
async def test_grep_keeps_an_oversized_matching_line_visible(tmp_path: Path) -> None:
    long_line = "x" * 130_000 + "needle" + "y" * 10_000
    (tmp_path / "huge-line.txt").write_text(long_line, encoding="utf-8")
    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)

    result = await tool.execute(
        pattern="needle",
        path="huge-line.txt",
        fixed_strings=True,
        context_before=0,
        context_after=0,
    )

    assert "huge-line.txt:1" in result
    assert "needle" in result
    assert "No matches found" not in result
    assert len(result) < GrepTool._MAX_RESULT_CHARS


@pytest.mark.asyncio
async def test_grep_size_limit_returns_a_resumable_offset(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    content = "\n".join(f"needle-{line}-" + "x" * 80 for line in range(1, 11))
    (tmp_path / "many.txt").write_text(content, encoding="utf-8")
    monkeypatch.setattr(GrepTool, "_MAX_RESULT_CHARS", 350)
    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)

    first = await tool.execute(
        pattern="needle",
        path="many.txt",
        fixed_strings=True,
        context_before=0,
        context_after=0,
        head_limit=10,
    )
    continuation = re.search(r"use offset=(\d+) to continue", first)

    assert continuation is not None
    next_offset = int(continuation.group(1))
    assert next_offset > 0

    second = await tool.execute(
        pattern="needle",
        path="many.txt",
        fixed_strings=True,
        context_before=0,
        context_after=0,
        head_limit=10,
        offset=next_offset,
    )
    first_headers = {
        line for line in first.splitlines() if line.startswith("many.txt:")
    }
    second_headers = {
        line for line in second.splitlines() if line.startswith("many.txt:")
    }
    assert second_headers
    assert first_headers.isdisjoint(second_headers)


@pytest.mark.asyncio
async def test_grep_reports_an_invalid_pdf_page_range(tmp_path: Path) -> None:
    from pypdf import PdfWriter

    pdf_path = tmp_path / "one-page.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with pdf_path.open("wb") as output:
        writer.write(output)

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(pattern="needle", path="one-page.pdf", pages="bad")

    assert result.startswith("Error: Invalid PDF page range 'bad'")
    assert "binary/unreadable" not in result

    out_of_bounds = await tool.execute(
        pattern="needle",
        path="one-page.pdf",
        pages="99",
    )
    assert out_of_bounds == (
        "Error: Invalid PDF page range '99': document has 1 page; "
        "use a page number or range within 1-1."
    )


@pytest.mark.asyncio
async def test_grep_supports_case_insensitive_search(tmp_path: Path) -> None:
    (tmp_path / "memory").mkdir()
    (tmp_path / "memory" / "HISTORY.md").write_text(
        "[2026-04-02 10:00] OAuth token rotated\n",
        encoding="utf-8",
    )

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="oauth",
        path="memory/HISTORY.md",
        case_insensitive=True,
        output_mode="content",
    )

    assert "memory/HISTORY.md:1" in result
    assert "OAuth token rotated" in result


@pytest.mark.asyncio
async def test_grep_type_filter_limits_files(tmp_path: Path) -> None:
    (tmp_path / "src").mkdir()
    (tmp_path / "src" / "a.py").write_text("needle\n", encoding="utf-8")
    (tmp_path / "src" / "b.md").write_text("needle\n", encoding="utf-8")

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="needle",
        path="src",
        type="py",
        output_mode="files_with_matches",
    )

    assert result.splitlines() == ["src/a.py"]


@pytest.mark.asyncio
async def test_grep_fixed_strings_treats_regex_chars_literally(tmp_path: Path) -> None:
    (tmp_path / "memory").mkdir()
    (tmp_path / "memory" / "HISTORY.md").write_text(
        "[2026-04-02 10:00] OAuth token rotated\n",
        encoding="utf-8",
    )

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="[2026-04-02 10:00]",
        path="memory/HISTORY.md",
        fixed_strings=True,
        output_mode="content",
    )

    assert "memory/HISTORY.md:1" in result
    assert "[2026-04-02 10:00] OAuth token rotated" in result


@pytest.mark.asyncio
async def test_grep_files_with_matches_mode_returns_unique_paths(tmp_path: Path) -> None:
    (tmp_path / "src").mkdir()
    a = tmp_path / "src" / "a.py"
    b = tmp_path / "src" / "b.py"
    a.write_text("needle\nneedle\n", encoding="utf-8")
    b.write_text("needle\n", encoding="utf-8")
    os.utime(a, (1, 1))
    os.utime(b, (2, 2))

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="needle",
        path="src",
        output_mode="files_with_matches",
    )

    assert result.splitlines() == ["src/b.py", "src/a.py"]


@pytest.mark.asyncio
async def test_grep_files_with_matches_supports_head_limit_and_offset(tmp_path: Path) -> None:
    (tmp_path / "src").mkdir()
    for name in ("a.py", "b.py", "c.py"):
        (tmp_path / "src" / name).write_text("needle\n", encoding="utf-8")

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="needle",
        path="src",
        output_mode="files_with_matches",
        head_limit=1,
        offset=1,
    )

    # Filesystem order is not deterministic across platforms, so just verify:
    # 1. Only one file path is returned (head_limit=1 after offset=1)
    # 2. The pagination info is correct
    assert "pagination: limit=1, offset=1" in result
    # Count non-empty lines that start with src/ (file paths)
    file_lines = [line for line in result.splitlines() if line.startswith("src/")]
    assert len(file_lines) == 1


@pytest.mark.asyncio
async def test_grep_count_mode_reports_counts_per_file(tmp_path: Path) -> None:
    (tmp_path / "logs").mkdir()
    (tmp_path / "logs" / "one.log").write_text("warn\nok\nwarn\n", encoding="utf-8")
    (tmp_path / "logs" / "two.log").write_text("warn\n", encoding="utf-8")

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="warn",
        path="logs",
        output_mode="count",
    )

    assert "logs/one.log: 2" in result
    assert "logs/two.log: 1" in result
    assert "total matches: 3 in 2 files" in result


@pytest.mark.asyncio
async def test_grep_files_with_matches_mode_respects_max_results(tmp_path: Path) -> None:
    (tmp_path / "src").mkdir()
    files = []
    for idx, name in enumerate(("a.py", "b.py", "c.py"), start=1):
        file_path = tmp_path / "src" / name
        file_path.write_text("needle\n", encoding="utf-8")
        os.utime(file_path, (idx, idx))
        files.append(file_path)

    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(
        pattern="needle",
        path="src",
        output_mode="files_with_matches",
        max_results=2,
    )

    assert result.splitlines()[:2] == ["src/c.py", "src/b.py"]
    assert "pagination: limit=2, offset=0" in result


@pytest.mark.asyncio
async def test_grep_reports_skipped_binary_and_large_files(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    (tmp_path / "binary.bin").write_bytes(b"\x00\x01\x02")
    (tmp_path / "large.txt").write_text("x" * 20, encoding="utf-8")

    monkeypatch.setattr(GrepTool, "_MAX_FILE_BYTES", 10)
    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    result = await tool.execute(pattern="needle", path=".")

    assert "No matches found" in result
    assert "skipped 1 binary/unreadable files" in result
    assert "skipped 1 large files" in result


@pytest.mark.asyncio
async def test_grep_uses_a_larger_bounded_limit_for_an_explicit_file(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    large_file = tmp_path / "history.jsonl"
    large_file.write_text("needle\n" + "x" * 20, encoding="utf-8")
    monkeypatch.setattr(GrepTool, "_MAX_FILE_BYTES", 10)
    monkeypatch.setattr(GrepTool, "_MAX_EXPLICIT_FILE_BYTES", 100)
    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)

    explicit_result = await tool.execute(
        pattern="needle",
        path=str(large_file),
        output_mode="content",
    )
    directory_result = await tool.execute(pattern="needle", path=".")
    monkeypatch.setattr(GrepTool, "_MAX_EXPLICIT_FILE_BYTES", 10)
    capped_result = await tool.execute(pattern="needle", path=str(large_file))

    assert "needle" in explicit_result
    assert "skipped 1 large files" in directory_result
    assert "skipped 1 large files" in capped_result


def test_grep_schema_is_concise_and_keeps_legacy_aliases_hidden(tmp_path: Path) -> None:
    tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)
    properties = tool.parameters["properties"]

    assert len(tool.description) < 150
    assert "2 MB" not in tool.description
    assert "100 MB" not in tool.description
    assert "head_limit" in properties
    assert "max_matches" not in properties
    assert "max_results" not in properties


@pytest.mark.asyncio
async def test_search_tools_reject_paths_outside_workspace(tmp_path: Path) -> None:
    outside = tmp_path.parent / "outside-search.txt"
    outside.write_text("secret\n", encoding="utf-8")

    grep_tool = GrepTool(workspace=tmp_path, allowed_dir=tmp_path)

    grep_result = await grep_tool.execute(pattern="secret", path=str(outside))

    assert grep_result.startswith("Error:")


def test_agent_loop_registers_grep(tmp_path: Path) -> None:
    bus = MessageBus()
    provider = MagicMock()
    provider.get_default_model.return_value = "test-model"

    loop = AgentLoop(bus=bus, provider=provider, workspace=tmp_path, model="test-model")

    assert "find_files" in loop.tools.tool_names
    assert "grep" in loop.tools.tool_names


@pytest.mark.asyncio
async def test_subagent_registers_grep(tmp_path: Path) -> None:
    bus = MessageBus()
    provider = MagicMock()
    provider.get_default_model.return_value = "test-model"
    provider.generation = GenerationSettings()
    mgr = SubagentManager(
        workspace=tmp_path,
        bus=bus,
        max_tool_result_chars=4096,
    )
    captured: dict[str, list[str]] = {}

    async def fake_run(spec):
        captured["tool_names"] = spec.tools.tool_names
        return SimpleNamespace(
            stop_reason="ok",
            final_content="done",
            tool_events=[],
            error=None,
        )

    mgr.runner.run = fake_run
    mgr._announce_result = AsyncMock()

    status = SubagentStatus(task_id="sub-1", label="label", task_description="search task", started_at=time.monotonic())
    await mgr._run_subagent(
        "sub-1",
        "search task",
        "label",
        {"channel": "cli", "chat_id": "direct"},
        status,
        LLMRuntime.capture(provider, "test-model", context_window_tokens=128_000),
    )

    assert "find_files" in captured["tool_names"]
    assert "grep" in captured["tool_names"]


def test_subagent_prompt_respects_disabled_skills(tmp_path: Path) -> None:
    bus = MessageBus()
    skills_dir = tmp_path / "skills"
    (skills_dir / "alpha").mkdir(parents=True)
    (skills_dir / "alpha" / "SKILL.md").write_text("# Alpha\n\nhidden\n", encoding="utf-8")
    (skills_dir / "beta").mkdir(parents=True)
    (skills_dir / "beta" / "SKILL.md").write_text("# Beta\n\nshown\n", encoding="utf-8")

    mgr = SubagentManager(
        workspace=tmp_path,
        bus=bus,
        max_tool_result_chars=4096,
        disabled_skills=["alpha"],
    )

    prompt = mgr._build_subagent_prompt()

    assert "alpha" not in prompt
    assert "beta" in prompt
