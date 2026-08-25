"""Document text extraction utilities for nanobot."""

import mimetypes
from collections.abc import Iterator
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from zipfile import BadZipFile, ZipFile

from loguru import logger

from nanobot.utils.helpers import detect_image_mime

_MAX_TEXT_LENGTH = 200_000
_MAX_EXTRACT_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
_MAX_OFFICE_ARCHIVE_MEMBERS = 10_000
_MAX_OFFICE_UNCOMPRESSED_SIZE = 256 * 1024 * 1024  # 256 MB
_MAX_OFFICE_MEMBER_SIZE = 128 * 1024 * 1024  # 128 MB
_MAX_DOCX_TABLE_CELLS = 100_000
_MAX_DOCX_TABLE_DEPTH = 8
_MAX_PDF_CONTENT_STREAM_SIZE = 32 * 1024 * 1024  # 32 MB per page
_MAX_PDF_ATTACHMENT_PAGES = 100


class _TextCollector:
    """Build bounded parser output without retaining the full document text."""

    def __init__(self, limit: int) -> None:
        self.limit = limit
        self.parts: list[str] = []
        self.length = 0
        self.truncated = False

    def add(self, text: str, *, separator: str = "") -> bool:
        if not text:
            return True
        prefix = separator if self.parts else ""
        chunk = prefix + text
        remaining = self.limit - self.length
        if len(chunk) > remaining:
            if remaining > 0:
                self.parts.append(chunk[:remaining])
                self.length += remaining
            self.truncated = True
            return False
        self.parts.append(chunk)
        self.length += len(chunk)
        return True

    def render(self) -> str:
        text = "".join(self.parts)
        if self.truncated:
            text += f"... (truncated at {self.limit} chars)"
        return text


class PdfSafetyError(Exception):
    """Raised when a PDF exceeds a parser safety boundary."""


class PdfPageRangeError(Exception):
    """Raised when a requested PDF page range is invalid."""


class DocxSafetyError(Exception):
    """Raised when a DOCX table exceeds a parser safety boundary."""


class DocumentExtractionError(Exception):
    """Raised when a document cannot be opened for incremental extraction."""


@dataclass(frozen=True, slots=True)
class PdfExtraction:
    text: str
    total_pages: int
    start_page: int
    end_page: int


@dataclass(frozen=True, slots=True)
class LocatedDocumentLine:
    """One searchable document line with a stable, human-readable locator."""

    text: str
    extracted_line: int
    locator: str
    searchable: bool = True


@dataclass(frozen=True, slots=True)
class DocumentLineSource:
    """Incremental document lines plus an optional next PDF page range."""

    lines: Iterator[LocatedDocumentLine]
    continuation: str | None = None


def extract_text(path: str | Path) -> str | None:
    """Extract text from a file.

    Args:
        path: Path to the file.

    Returns:
        Extracted text as string, None for unsupported types,
        or error string for failures.
    """
    path = Path(path)
    if error := _extraction_path_error(path):
        return error

    ext = path.suffix.lower()

    # Parsers stay lazy even though they are bundled so idle processes do not
    # retain their import cost (see issue #3422).
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext == ".docx":
        return _extract_docx(path)
    elif ext == ".xlsx":
        return _extract_xlsx(path)
    elif ext == ".pptx":
        return _extract_pptx(path)
    elif _is_text_extension(ext):
        return _extract_text_file(path)
    elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        # Image files - for future OCR support
        return f"[image: {path.name}]"
    else:
        # Unsupported extension
        return None


def open_document_line_source(
    path: str | Path,
    *,
    pages: str | None = None,
) -> DocumentLineSource | None:
    """Open a document as an incremental stream of extracted lines.

    Unlike :func:`extract_text`, this interface does not apply the attachment
    text preview limit. Parser/file safety limits still apply. Lines that are
    useful only for the rendered document view (for example sheet headers and
    blank separators) have ``searchable=False`` so range reads can retain them
    without making grep match synthetic text.
    """
    path = Path(path)
    ext = path.suffix.lower()
    if ext not in {".pdf", ".docx", ".xlsx", ".pptx"}:
        return None
    if error := _extraction_path_error(path):
        raise DocumentExtractionError(_clean_extraction_error(error))
    if ext == ".pdf":
        return _open_pdf_line_source(path, pages)
    if ext == ".docx":
        return _open_docx_line_source(path)
    if ext == ".xlsx":
        return _open_xlsx_line_source(path)
    return _open_pptx_line_source(path)


def _clean_extraction_error(error: str) -> str:
    if error.startswith("[error:") and error.endswith("]"):
        return error[len("[error:") : -1].strip()
    return error


def _check_office_archive(path: Path) -> None:
    if error := _office_archive_error(path):
        raise DocumentExtractionError(_clean_extraction_error(error))


def _open_pdf_line_source(path: Path, pages: str | None) -> DocumentLineSource:
    try:
        from pypdf import PdfReader

        reader = PdfReader(path, strict=False)
        total_pages = len(reader.pages)
        if total_pages == 0:
            return DocumentLineSource(iter(()))
        start, requested_end = _parse_pdf_page_range(pages, total_pages)
    except PdfPageRangeError:
        raise
    except Exception as e:
        raise DocumentExtractionError(f"failed to open PDF: {e!s}") from e

    end = min(requested_end, start + _MAX_PDF_ATTACHMENT_PAGES - 1)
    continuation = None
    if end < total_pages - 1:
        next_start = end + 2
        next_end = min(end + 1 + _MAX_PDF_ATTACHMENT_PAGES, total_pages)
        continuation = f"pages='{next_start}-{next_end}'"

    def iter_lines() -> Iterator[LocatedDocumentLine]:
        extracted_line = 0
        wrote_page = False
        for index in range(start, end + 1):
            page = reader.pages[index]
            contents = page.get_contents()
            if contents is not None:
                stream_size = len(contents.get_data())
                if stream_size > _MAX_PDF_CONTENT_STREAM_SIZE:
                    raise PdfSafetyError(
                        f"page {index + 1} content stream exceeds "
                        f"{_MAX_PDF_CONTENT_STREAM_SIZE // (1024 * 1024)} MB limit"
                    )
            text = (page.extract_text() or "").strip()
            if not text:
                continue
            if wrote_page:
                extracted_line += 1
                yield LocatedDocumentLine("", extracted_line, "", searchable=False)
            extracted_line += 1
            yield LocatedDocumentLine(
                f"--- Page {index + 1} ---",
                extracted_line,
                "",
                searchable=False,
            )
            page_line = 0
            for text_line in text.splitlines():
                extracted_line += 1
                if not text_line:
                    yield LocatedDocumentLine("", extracted_line, "", searchable=False)
                    continue
                page_line += 1
                yield LocatedDocumentLine(
                    text_line,
                    extracted_line,
                    f"page={index + 1},line={page_line}",
                )
            wrote_page = True

    return DocumentLineSource(iter_lines(), continuation=continuation)


def _open_xlsx_line_source(path: Path) -> DocumentLineSource:
    _check_office_archive(path)
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise DocumentExtractionError("openpyxl not installed") from e
    try:
        workbook = load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        raise DocumentExtractionError(f"failed to open XLSX: {e!s}") from e

    def iter_lines() -> Iterator[LocatedDocumentLine]:
        extracted_line = 0
        wrote_document_content = False
        try:
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                wrote_header = False
                for row_index, row in enumerate(worksheet.iter_rows(values_only=True), 1):
                    row_text = "\t".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if not row_text.strip():
                        continue
                    if not wrote_header:
                        if wrote_document_content:
                            extracted_line += 1
                            yield LocatedDocumentLine(
                                "", extracted_line, "", searchable=False
                            )
                        extracted_line += 1
                        yield LocatedDocumentLine(
                            f"--- Sheet: {sheet_name} ---",
                            extracted_line,
                            "",
                            searchable=False,
                        )
                        wrote_header = True
                        wrote_document_content = True
                    extracted_line += 1
                    yield LocatedDocumentLine(
                        row_text,
                        extracted_line,
                        f"sheet={sheet_name!r},row={row_index}",
                    )
        finally:
            workbook.close()

    return DocumentLineSource(iter_lines())


def _open_pptx_line_source(path: Path) -> DocumentLineSource:
    _check_office_archive(path)
    try:
        from pptx import Presentation as PptxPresentation
    except ImportError as e:
        raise DocumentExtractionError("python-pptx not installed") from e
    try:
        presentation = PptxPresentation(str(path))
    except Exception as e:
        raise DocumentExtractionError(f"failed to open PPTX: {e!s}") from e

    def iter_lines() -> Iterator[LocatedDocumentLine]:
        extracted_line = 0
        wrote_slide = False
        for slide_number, slide in enumerate(presentation.slides, 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                _collect_pptx_shape_text(shape, slide_text)
            rendered_lines = [line for text in slide_text for line in text.splitlines()]
            if not rendered_lines:
                continue
            if wrote_slide:
                extracted_line += 1
                yield LocatedDocumentLine("", extracted_line, "", searchable=False)
            extracted_line += 1
            yield LocatedDocumentLine(
                f"--- Slide {slide_number} ---",
                extracted_line,
                "",
                searchable=False,
            )
            slide_line = 0
            for text_line in rendered_lines:
                extracted_line += 1
                if not text_line:
                    yield LocatedDocumentLine("", extracted_line, "", searchable=False)
                    continue
                slide_line += 1
                yield LocatedDocumentLine(
                    text_line,
                    extracted_line,
                    f"slide={slide_number},line={slide_line}",
                )
            wrote_slide = True

    return DocumentLineSource(iter_lines())


def _open_docx_line_source(path: Path) -> DocumentLineSource:
    _check_office_archive(path)
    try:
        from docx import Document as DocxDocument
        from docx.table import Table, _Cell  # pyright: ignore[reportPrivateUsage]
        from docx.text.paragraph import Paragraph
    except ImportError as e:
        raise DocumentExtractionError("python-docx not installed") from e
    try:
        document = DocxDocument(str(path))
    except Exception as e:
        raise DocumentExtractionError(f"failed to open DOCX: {e!s}") from e

    def iter_lines() -> Iterator[LocatedDocumentLine]:
        table_cell_count = 0

        def cell_text(cell: _Cell, depth: int) -> str:
            parts: list[str] = []
            for block in cell.iter_inner_content():
                if isinstance(block, Paragraph):
                    text = " ".join(block.text.split())
                    if text:
                        parts.append(text)
                elif isinstance(block, Table):  # pyright: ignore[reportUnnecessaryIsInstance]
                    parts.extend(
                        row.replace("\t", " | ") for row in table_rows(block, depth + 1)
                    )
            return " ".join(parts)

        def table_rows(table: Table, depth: int) -> Iterator[str]:
            nonlocal table_cell_count
            if depth > _MAX_DOCX_TABLE_DEPTH:
                raise DocxSafetyError(
                    f"table nesting exceeds {_MAX_DOCX_TABLE_DEPTH} levels"
                )
            for row in table.rows:
                cells: list[str] = []
                for tc in row._tr.tc_lst:  # pyright: ignore[reportPrivateUsage]
                    table_cell_count += 1
                    if table_cell_count > _MAX_DOCX_TABLE_CELLS:
                        raise DocxSafetyError(
                            f"document contains more than {_MAX_DOCX_TABLE_CELLS} table cells"
                        )
                    cells.append(cell_text(_Cell(tc, table), depth))
                if any(cells):
                    yield "\t".join(cells)

        def blocks() -> Iterator[tuple[str, bool]]:
            for block in document.iter_inner_content():
                if isinstance(block, Paragraph):
                    text = block.text.strip()
                    if text:
                        yield text, True
                    continue
                if not isinstance(block, Table):  # pyright: ignore[reportUnnecessaryIsInstance]
                    continue
                first_row = True
                for row_text in table_rows(block, 1):
                    yield row_text, first_row
                    first_row = False

        extracted_line = 0
        paragraph = 0
        wrote_content = False
        for text, separate in blocks():
            if wrote_content and separate:
                extracted_line += 1
                yield LocatedDocumentLine("", extracted_line, "", searchable=False)
            for text_line in text.splitlines():
                extracted_line += 1
                if not text_line:
                    yield LocatedDocumentLine("", extracted_line, "", searchable=False)
                    continue
                paragraph += 1
                yield LocatedDocumentLine(
                    text_line,
                    extracted_line,
                    f"paragraph={paragraph}",
                )
            wrote_content = True

    return DocumentLineSource(iter_lines())


def _extraction_path_error(path: Path) -> str | None:
    if not path.exists():
        return f"[error: file not found: {path}]"
    try:
        if path.stat().st_size > _MAX_EXTRACT_FILE_SIZE:
            return f"[error: file exceeds {_MAX_EXTRACT_FILE_SIZE // (1024 * 1024)} MB limit]"
    except OSError as e:
        return f"[error: failed to inspect file: {e!s}]"
    return None


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using pypdf."""
    try:
        result = extract_pdf_pages(
            path,
            max_pages=_MAX_PDF_ATTACHMENT_PAGES,
            max_chars=_MAX_TEXT_LENGTH,
        )
        text = result.text
        if result.end_page < result.total_pages - 1:
            text += f"\n\n(Showing pages 1-{result.end_page + 1} of {result.total_pages}.)"
        return text
    except Exception as e:
        logger.exception("Failed to extract PDF {}", path)
        return f"[error: failed to extract PDF: {e!s}]"


def extract_pdf_pages(
    path: Path,
    *,
    pages: str | None = None,
    max_pages: int = _MAX_PDF_ATTACHMENT_PAGES,
    max_chars: int = _MAX_TEXT_LENGTH,
) -> PdfExtraction:
    """Extract a bounded PDF page range using the bundled pypdf reader."""
    from pypdf import PdfReader

    reader = PdfReader(path, strict=False)
    total_pages = len(reader.pages)
    if total_pages == 0:
        return PdfExtraction("", 0, 0, -1)

    start, end = _parse_pdf_page_range(pages, total_pages)
    end = min(end, start + max_pages - 1)
    collector = _TextCollector(max_chars)
    for index in range(start, end + 1):
        page = reader.pages[index]
        contents = page.get_contents()
        if contents is not None:
            stream_size = len(contents.get_data())
            if stream_size > _MAX_PDF_CONTENT_STREAM_SIZE:
                raise PdfSafetyError(
                    f"page {index + 1} content stream exceeds "
                    f"{_MAX_PDF_CONTENT_STREAM_SIZE // (1024 * 1024)} MB limit"
                )
        text = (page.extract_text() or "").strip()
        if text and not collector.add(f"--- Page {index + 1} ---\n{text}", separator="\n\n"):
            end = index
            break
    return PdfExtraction(collector.render(), total_pages, start, end)


def _parse_pdf_page_range(pages: str | None, total_pages: int) -> tuple[int, int]:
    if not pages:
        return 0, total_pages - 1
    page_word = "page" if total_pages == 1 else "pages"
    guidance = (
        f"document has {total_pages} {page_word}; "
        f"use a page number or range within 1-{total_pages}"
    )
    values = pages.strip().split("-")
    if len(values) not in {1, 2}:
        raise PdfPageRangeError(guidance)
    try:
        start = int(values[0])
        end = int(values[-1])
    except ValueError as e:
        raise PdfPageRangeError(guidance) from e
    if start < 1 or end < start or start > total_pages:
        raise PdfPageRangeError(guidance)
    return start - 1, min(end, total_pages) - 1


def _render_document_preview(source: DocumentLineSource) -> str:
    """Render a bounded attachment preview from the canonical line stream."""
    collector = _TextCollector(_MAX_TEXT_LENGTH)
    iterator = source.lines
    first_line = True
    try:
        for line in iterator:
            if not first_line and not collector.add("\n"):
                break
            first_line = False
            if line.text and not collector.add(line.text):
                break
        return collector.render()
    finally:
        close = getattr(iterator, "close", None)
        if close is not None:
            close()


def _extract_docx(path: Path) -> str:
    """Extract a bounded DOCX attachment preview."""
    try:
        return _render_document_preview(_open_docx_line_source(path))
    except DocxSafetyError as e:
        return f"[error: unsafe DOCX: {e!s}]"
    except DocumentExtractionError as e:
        return f"[error: {e!s}]"
    except Exception as e:
        logger.exception("Failed to extract DOCX {}", path)
        return f"[error: failed to extract DOCX: {e!s}]"


def _extract_xlsx(path: Path) -> str:
    """Extract a bounded XLSX attachment preview."""
    try:
        return _render_document_preview(_open_xlsx_line_source(path))
    except DocumentExtractionError as e:
        return f"[error: {e!s}]"
    except Exception as e:
        logger.exception("Failed to extract XLSX {}", path)
        return f"[error: failed to extract XLSX: {e!s}]"


def _extract_pptx(path: Path) -> str:
    """Extract a bounded PPTX attachment preview."""
    try:
        return _render_document_preview(_open_pptx_line_source(path))
    except DocumentExtractionError as e:
        return f"[error: {e!s}]"
    except Exception as e:
        logger.exception("Failed to extract PPTX {}", path)
        return f"[error: failed to extract PPTX: {e!s}]"


def _collect_pptx_shape_text(shape: Any, out: list[str]) -> None:
    """Collect text from a PPTX shape, recursing into groups and tables.

    Groups have ``has_text_frame=False`` and must be walked via ``.shapes``;
    tables are GraphicFrame objects whose cell text lives under ``.table``.
    """
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub in sub_shapes:
            _collect_pptx_shape_text(sub, out)
        return

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cell for cell in cells if cell)
            if line:
                out.append(line)
        return

    text = getattr(shape, "text", "")
    if text:
        out.append(text)


def _office_archive_error(path: Path) -> str | None:
    """Reject oversized or encrypted OOXML containers before parsing XML."""
    try:
        with ZipFile(path) as archive:
            members = archive.infolist()
    except (BadZipFile, OSError) as e:
        return f"[error: invalid Office document: {e!s}]"
    if len(members) > _MAX_OFFICE_ARCHIVE_MEMBERS:
        return f"[error: Office document contains too many files ({len(members)})]"
    total_size = 0
    for member in members:
        if member.flag_bits & 0x1:
            return "[error: encrypted Office documents are not supported]"
        if member.file_size > _MAX_OFFICE_MEMBER_SIZE:
            return "[error: Office document contains an oversized internal file]"
        total_size += member.file_size
        if total_size > _MAX_OFFICE_UNCOMPRESSED_SIZE:
            limit_mb = _MAX_OFFICE_UNCOMPRESSED_SIZE / (1024 * 1024)
            return f"[error: Office document expands beyond the {limit_mb:g} MB safety limit]"
    return None


def _extract_text_file(path: Path) -> str:
    """Extract text from a plain text file."""
    try:
        # Try UTF-8 first, then latin-1 fallback
        try:
            content = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            content = path.read_text(encoding="latin-1")
        return _truncate(content, _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.exception("Failed to read text file {}", path)
        return f"[error: failed to read file: {e!s}]"


def _truncate(text: str, max_length: int) -> str:
    """Truncate text with a suffix indicating truncation."""
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"... (truncated, {len(text)} chars total)"


def _is_text_extension(ext: str) -> bool:
    """Check if extension is a text format."""
    return ext in {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".log",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".cfg",
    }


# ---------------------------------------------------------------------------
# High-level helper: split images from on-demand attachment references
# ---------------------------------------------------------------------------


def is_image_file(path: str) -> bool:
    """Check whether *path* looks like an image file.

    Uses magic-byte detection (reads first 16 bytes) with a ``mimetypes``
    extension-based fallback.
    """
    p = Path(path)
    mime: str | None = None
    if p.is_file():
        try:
            with p.open("rb") as f:
                mime = detect_image_mime(f.read(16))
        except OSError:
            mime = None
    if not mime:
        mime = mimetypes.guess_type(path)[0]
    return bool(mime and mime.startswith("image/"))


def _canonical_local_media_path(path: str) -> str:
    """Return an existing local media file as an absolute path."""
    try:
        candidate = Path(path).expanduser()
        if candidate.is_file():
            return str(candidate.resolve(strict=False))
    except (OSError, RuntimeError, TypeError, ValueError):
        pass
    return path


def reference_non_image_attachments(
    content: str, media: list[str],
) -> tuple[str, list[str]]:
    """Reference non-image attachments without reading file content.

    Image paths are preserved for downstream vision-block construction.
    Non-image paths are appended as ``[Attachment: path]`` references so the
    model can inspect them on demand with ``read_file`` or pass the original
    path to another tool that needs exact file bytes.
    """
    image_paths: list[str] = []
    attachment_refs: list[str] = []
    for path in media:
        path = _canonical_local_media_path(path)
        if is_image_file(path):
            image_paths.append(path)
        else:
            attachment_refs.append(f"[Attachment: {path}]")
    if attachment_refs:
        suffix = "\n".join(attachment_refs)
        content = f"{content}\n\n{suffix}" if content else suffix
    return content, image_paths
